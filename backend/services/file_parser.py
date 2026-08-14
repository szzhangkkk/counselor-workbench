import os
import re
import io
from typing import List, Dict, Any

import pandas as pd
from docx import Document
from pptx import Presentation

# ---------------------------------------------------------------------------
# 台账字段 -> 可能的 Excel/表头别名（中文智能映射）
# ---------------------------------------------------------------------------
FIELD_ALIASES: Dict[str, Dict[str, List[str]]] = {
    "students": {
        "student_id": ["学号", "学生学号", "学号(必须)", "ID"],
        "name": ["姓名", "学生姓名"],
        "gender": ["性别"],
        "major": ["专业", "专业名称"],
        "class_name": ["班级", "行政班级", "班别"],
        "grade": ["年级", "届", "届别", "级", "入学年份", "入学年度"],
        "phone": ["手机号", "手机号码", "联系电话", "电话"],
        "email": ["邮箱", "电子邮件", "Email"],
        "dormitory": ["宿舍", "寝室", "宿舍号"],
        "political_status": ["政治面貌"],
        "hometown": ["生源地", "籍贯", "家乡", "家庭所在地"],
        "birth_date": ["出生日期", "生日"],
        "guardian_name": ["监护人", "家长姓名"],
        "guardian_phone": ["监护人电话", "家长电话"],
        "special_info": ["备注", "特殊说明", "特殊信息", "其他"],
    },
    "employment": {
        "student_id": ["学号"],
        "name": ["姓名"],
        "class_name": ["班级"],
        "status": ["就业状态", "状态", "毕业去向", "就业去向"],
        "company": ["单位", "公司", "用人单位", "就业单位"],
        "position": ["职位", "岗位"],
        "salary": ["薪资", "月薪", "工资"],
        "offer_date": ["签约日期", "offer日期", "就职日期", "日期"],
        "contract_type": ["合同类型", "签约类型"],
        "region": ["地区", "单位所在地", "城市"],
        "remarks": ["备注", "说明"],
    },
    "psychology": {
        "student_id": ["学号"],
        "name": ["姓名"],
        "class_name": ["班级"],
        "assess_date": ["评估日期", "测评日期", "日期", "访谈日期"],
        "level": ["等级", "心理等级", "状态", "评估等级", "预警等级"],
        "category": ["类别", "问题类别", "类型"],
        "symptoms": ["表现", "症状", "问题描述", "情况"],
        "intervention": ["干预措施", "干预方案", "采取措施", "处理方式"],
        "counselor": ["咨询师", "辅导员", "负责人"],
        "next_follow_up": ["下次跟进", "回访日期", "下次回访"],
        "remarks": ["备注", "说明"],
    },
    "talks": {
        "student_id": ["学号"],
        "name": ["姓名"],
        "class_name": ["班级"],
        "talk_date": ["谈话日期", "日期", "时间"],
        "talk_type": ["谈话类型", "类型"],
        "topic": ["主题", "谈话主题", "话题"],
        "content": ["内容", "谈话内容", "详情"],
        "conclusion": ["结论", "谈话结论", "效果"],
        "counselor": ["谈话人", "辅导员", "记录人"],
    },
    "grades": {
        "student_id": ["学号"],
        "name": ["姓名"],
        "class_name": ["班级"],
        "semester": ["学期"],
        "course_name": ["课程", "课程名称"],
        "course_code": ["课程代码", "课程编号"],
        "credits": ["学分"],
        "score": ["成绩", "分数", "课程成绩"],
        "grade_point": ["绩点"],
        "rank": ["排名"],
        "make_up_score": ["补考成绩", "重修成绩"],
    },
    "attendance": {
        "student_id": ["学号"],
        "name": ["姓名"],
        "class_name": ["班级"],
        "course_name": ["课程", "课程名称"],
        "date": ["日期"],
        "period": ["节次", "时间段", "课时"],
        "status": ["状态", "出勤情况", "考勤"],
        "reason": ["原因", "请假原因"],
        "recorder": ["记录人", "考勤人"],
    },
}

# 每条记录必需的字段（缺少则跳过该行）
REQUIRED_FIELDS = {
    "students": ["name"],
    "employment": ["name"],
    "psychology": ["name"],
    "talks": ["name"],
    "grades": ["name", "course_name"],
    "attendance": ["name"],
}


class FileParserError(Exception):
    pass


def normalize_header(h: str) -> str:
    if h is None:
        return ""
    s = str(h).strip().replace("\n", "").replace(" ", "")
    # 全角 -> 半角
    out = []
    for ch in s:
        code = ord(ch)
        if code == 0x3000:
            code = 32
        elif 0xFF01 <= code <= 0xFF5E:
            code -= 0xFEE0
        out.append(chr(code))
    return "".join(out).lower()


def _all_aliases(ledger_type: str) -> set:
    return {
        normalize_header(n)
        for names in FIELD_ALIASES.get(ledger_type, {}).values()
        for n in names
    }


def build_header_map(columns: List[str], ledger_type: str) -> Dict[str, int]:
    """返回 {model_field: column_index} 映射。
    第一轮：精确匹配；第二轮：正向包含匹配（列名包含别名，如"学生姓名"包含"姓名"）。
    每列最多被使用一次，避免"姓名"被"家长姓名"反向误配。"""
    aliases = FIELD_ALIASES.get(ledger_type, {})
    header_map: Dict[str, int] = {}
    norm_cols = [(i, normalize_header(c)) for i, c in enumerate(columns)]
    used: set = set()

    # 第一轮：精确匹配
    for field, names in aliases.items():
        for name in names:
            key = normalize_header(name)
            if not key:
                continue
            for idx, nc in norm_cols:
                if idx not in used and nc == key:
                    header_map[field] = idx
                    used.add(idx)
                    break
            if field in header_map:
                break

    # 第二轮：正向包含匹配，取最长别名
    for field, names in aliases.items():
        if field in header_map:
            continue
        best_idx, best_len = None, 0
        for name in names:
            key = normalize_header(name)
            if not key or len(key) < 2:
                continue
            for idx, nc in norm_cols:
                if idx in used or not nc:
                    continue
                if key in nc and len(key) > best_len:
                    best_idx, best_len = idx, len(key)
        if best_idx is not None:
            header_map[field] = best_idx
            used.add(best_idx)
    return header_map


def _score_row(columns: List[str], ledger_type: str) -> int:
    """评估一行作为表头的匹配度"""
    alias_set = _all_aliases(ledger_type)
    return sum(1 for c in columns if normalize_header(c) in alias_set)


def find_header_row(rows: List[List[str]], ledger_type: str) -> tuple:
    """自动定位表头行（扫描前 10 行），返回 (行索引, 匹配分数)"""
    best_idx, best_score = 0, -1
    for i, row in enumerate(rows[:10]):
        score = _score_row(row, ledger_type)
        if score > best_score:
            best_idx, best_score = i, score
    return best_idx, best_score


def _clean_value(v: Any) -> Any:
    if v is None:
        return ""
    if isinstance(v, float) and v.is_integer():
        return str(int(v))
    if isinstance(v, pd.Timestamp):
        return v.strftime("%Y-%m-%d")
    return str(v).strip()


def _row_to_dict(row, header_map: Dict[str, int], ledger_type: str) -> Dict[str, Any]:
    out: Dict[str, Any] = {}
    for field, idx in header_map.items():
        if idx < len(row):
            out[field] = _clean_value(row[idx])
    return out


def _row_valid(data: Dict[str, Any], ledger_type: str) -> bool:
    for f in REQUIRED_FIELDS.get(ledger_type, []):
        if not data.get(f):
            return False
    return True


# ---------------------------------------------------------------------------
# 各格式解析 -> List[List[str]] 表格行
# ---------------------------------------------------------------------------
def _parse_excel(content: bytes, filename: str) -> List[List[str]]:
    try:
        df = pd.read_excel(io.BytesIO(content), header=None)
    except Exception:
        df = pd.read_excel(io.BytesIO(content), header=None, engine="openpyxl")
    if df.empty:
        return []
    return df.astype(object).where(df.notna(), None).values.tolist()


def _parse_markdown(content: str) -> List[List[str]]:
    """解析 markdown 中第一个表格为二维数组"""
    lines = content.splitlines()
    table_lines = []
    in_table = False
    for ln in lines:
        s = ln.strip()
        if s.startswith("|"):
            in_table = True
            table_lines.append(s)
        elif in_table:
            break
    if not table_lines:
        return []
    rows = []
    for i, ln in enumerate(table_lines):
        cells = [c.strip() for c in ln.strip().strip("|").split("|")]
        # 跳过分隔行 |---|---|
        if all(re.fullmatch(r":?-{2,}:?", c) for c in cells if c):
            continue
        rows.append(cells)
    return rows


def _parse_latex(content: str) -> List[List[str]]:
    """解析 LaTeX 中第一个 tabular 环境"""
    m = re.search(r"\\begin\{tabular\}.*?\}(.*?)\\end\{tabular\}", content, re.S)
    if not m:
        return []
    body = m.group(1)
    rows = []
    for line in body.split("\\\\"):
        line = re.sub(r"\\\\(?!\s)", "", line)
        if "\\hline" in line or not line.strip():
            continue
        cells = [re.sub(r"[\\{}]", "", c).strip() for c in line.split("&")]
        rows.append(cells)
    return rows


def _parse_pdf(content: bytes) -> List[List[str]]:
    import pdfplumber

    with pdfplumber.open(io.BytesIO(content)) as pdf:
        rows: List[List[str]] = []
        for page in pdf.pages:
            table = page.extract_table()
            if table:
                for r in table:
                    rows.append([(c or "").strip() for c in r])
            else:
                text = page.extract_text()
                if text:
                    for ln in text.splitlines():
                        rows.append([ln.strip()])
        return rows


def _parse_docx(content: bytes) -> List[List[str]]:
    doc = Document(io.BytesIO(content))
    rows: List[List[str]] = []
    for table in doc.tables:
        for r in table.rows:
            rows.append([c.text.strip() for c in r.cells])
    if not rows:
        for p in doc.paragraphs:
            if p.text.strip():
                rows.append([p.text.strip()])
    return rows


def _parse_pptx(content: bytes) -> List[List[str]]:
    prs = Presentation(io.BytesIO(content))
    rows: List[List[str]] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for r in shape.table.rows:
                    rows.append([c.text.strip() for c in r.cells])
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    rows.append([text])
    return rows


def parse_file(content: bytes, filename: str, ledger_type: str) -> List[Dict[str, Any]]:
    """解析文件并返回结构化记录列表"""
    ext = os.path.splitext(filename)[1].lower()
    text = ""
    rows: List[List[str]] = []

    if ext in (".xlsx", ".xls"):
        rows = _parse_excel(content, filename)
    elif ext == ".md" or ext == ".markdown":
        text = content.decode("utf-8", errors="ignore")
        rows = _parse_markdown(text)
    elif ext == ".tex":
        text = content.decode("utf-8", errors="ignore")
        rows = _parse_latex(text)
    elif ext == ".pdf":
        rows = _parse_pdf(content)
    elif ext == ".docx":
        rows = _parse_docx(content)
    elif ext == ".pptx":
        rows = _parse_pptx(content)
    else:
        raise FileParserError(f"不支持的文件类型: {ext}")

    if not rows:
        raise FileParserError("文件中未解析到任何表格或数据")

    # 自动定位表头行（标题行/汇总行会被跳过）
    header_idx, score = find_header_row(rows, ledger_type)
    header = [str(c) for c in rows[header_idx]]
    header_map = build_header_map(header, ledger_type)

    if score == 0:
        sample = rows[:2]
        preview = " / ".join(str(c) for c in sample[0]) if sample else ""
        raise FileParserError(
            f"无法识别表头列名（匹配分数 0）。\n"
            f"文件首行内容：{preview}\n"
            f"请确保首行为列标题，且包含可识别列名，如：学号、姓名、班级、"
            + "、".join(list(FIELD_ALIASES.get(ledger_type, {}).keys())[:6])
            + " 等。"
        )

    data_rows = rows[header_idx + 1 :]

    if not header_map:
        raise FileParserError(
            f"表头行被识别，但未匹配到可用字段。表头：{' / '.join(header)}"
        )

    records = []
    for row in data_rows:
        data = _row_to_dict(row, header_map, ledger_type)
        if _row_valid(data, ledger_type):
            records.append(data)

    return records
